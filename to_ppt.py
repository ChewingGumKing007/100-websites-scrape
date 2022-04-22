# -*- coding: utf-8 -*-
"""
Created on Wed May 19 13:08:59 2021

@author: Oscar Jiménez
"""
################### EXPORT TO PPT #####################

# pip install python-pptx
import sys
sys.path.append('../../libraries')

# from iota_utils.iota_utils import US_state_abrev
from libraries.data_tag import data_tag

from pptx import Presentation
from pptx.util import Inches
import datetime
import math
from pptx.util import Cm, Pt
from PIL import Image
#from PIL.Image import core as image
from pptx.dml.color import RGBColor
import pandas as pd
import re
import os
import comtypes.client
import numpy as np
from pptx.enum.text import MSO_VERTICAL_ANCHOR
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import PP_ALIGN



def to_ppt(datas, name, data_frame_name, data_dictionary,
           # df
           ):  
   
    """
    This function creates a ppt from a dataframe.
    :params datas: this an array that contains the following:
                    1. agenda: the agenda for creating the table content
                    2. name_graph: array with the name of the graphs
                    3. loc_graph: location of the graphs
                    4. source: string with the source of the dataset
                    5. fake_agenda: a fake agenda that list every each item of the agenda, even the ones which are not listed.
                    6. level: a integer that indicates the level in which the hyperlinks are going to be setted up
                    7. space_ocuppers: a array that indicates what elements of the fake agenda occupes exactly one slide.
                    8. methodology: array that contains all the information for the methodology slides.
    :params name: String that indicates the name of the script
    :params data_frame_name: array that contains keywords, answers and metadata
    :params data_dictionary: array that contains metadata
    :params df: the dataframe that contains all the values used to create the graphs.
    """
    
    # template
    watermark=datas[10]
    water_mark="SamplesPNG.png"
    
    #EXTRA_CONTENT=3 #THIS VARIABLE INDICATES HOW MANY SLIDES NEEDS TO BE COUNTED FROM THE BEGGING TO THE END IN 
                    #ORDER TO START THE COUNTER FOR THE HYPERLINKS
                    #THIS IS CALCULATED LIKE: THE NUMBER OF SLIDES BEFORE THE FIRST HYPERLINK NEEDS TO BE SET UP
                    #WITHOUT COUNTING THE LENGHT OF THE TABLE OF CONTENT.
                    #BECAUSE THE LENGTH IS CALCULATED LIKE: LEN(AGENDA_PG)+EXTRA_CONTENT
    
    
    prs = Presentation(r'ppt/Template.pptx')             
    
    #copy
    prs.save(r'reports/{}.pptx'.format(name))   

    """
    for i in df.columns:
        if df[i].dtype == np.float64 or df[i].dtype == np.int64:
            max_values=[[max(df[i]),df[i].argmax(),i]]
        #elif df[i].dtype == "String":
            #max_values=[[df[i].value_counts().idxmax(),i]]

    print(max_values)

    """
    
    
# =============================================================================
#  Introduction template
# =============================================================================
    """
    In this section we just append all the required information in the cover of the pptx
    with a single param called data_frame_name.
    :data_frame_name: comes directly from t_ppt function, it stores the keywords values and answers values
    """
    first_slide_template = prs.slide_layouts[0] 
    #keywords_values= data_tag.tag(data_frame_name)[0]
    #ds = "Keywords: "+", ".join( key.replace("#","") for key in keywords_values)
    #answers = data_tag.tag(data_frame_name)[1]
    #answers = ["-"+answer for answer in answers]
    #answers = "Abstract: "+", ".join(answers)
    
    
    first_slide = prs.slides.add_slide(first_slide_template) 
    #for shape in first_slide.placeholders:
    #    print('%d %s' % (shape.placeholder_format.idx, shape.name))
    
    sec=first_slide.shapes.placeholders[0].text_frame.paragraphs[0]
    #sec.clear()
    sec.top = Inches(1.89)
    sec.left = Inches(0.56)
    sec.width = Inches(11.28)
    sec.height = Inches(1.71)
    
    secon=sec.add_run()
    #secon.text = "k234454454"
    second=secon.font
    second.color.rgb = RGBColor(0, 0, 0)
    second.size=Pt(26)
    second.name="Mulish"
    second.bold=True
    #second.line_spacing=1.1
    
    secon.text = "{}".format(name)+"\nData set: "+data_frame_name
    
    first_slide.placeholders[1].top = Inches(6.57)
    first_slide.placeholders[1].left = Inches(0.56)
    first_slide.placeholders[1].width = Inches(11.28)
    first_slide.placeholders[1].height = Inches(0.48)
    first_slide.placeholders[1].text = "Last updated "+ str(datetime.date.today())#+"\n "+keywords +"\n "+answers   
    #above are tghe contents of the title slide...the date it was last updated and the title of the project.
    
    #saving the contents of "answers" to a variable
    answers = data_tag.tag(data_frame_name)[1]
    answers = ["-"+answer for answer in answers]
    answers = ""+". ".join(answers)
    answers=answers.split('-')[1:]
    paragraph_strs =answers
    #for shape in slide.placeholders:
    #    print('%d %s' % (shape.placeholder_format.idx, shape.name))
    
    
    #specify the dimensions of the answers' contents placeholder
    first_slide.shapes.placeholders[10].top = Inches(3.75)
    first_slide.shapes.placeholders[10].left = Inches(0.56)
    first_slide.shapes.placeholders[10].width = Inches(11.03)
    first_slide.shapes.placeholders[10].height = Inches(2.82)
    
    #adding a textframe to the contents of "abstract"
    fiki = first_slide.placeholders[10].text_frame#.add_paragraph()
    
    fiki.clear()  # remove any existing paragraphs, leaving one empty one
    
    
    #create first bullet point of the contents of "abstract" and specifying its specifications
    # parii = fiki.paragraphs[0]
    # parii.text = paragraph_strs[0]
    # parii.font.name="Mulish"
    # parii.font.size=Pt(19)
    # parii.line_spacing=1.1
    # parii.bold=True
    
    #creating other bullet points of the "abstract" contents
    for para_str in paragraph_strs:
        parii = fiki.add_paragraph()
        parii.text = para_str
        parii.font.name="Mulish"
        parii.font.size=Pt(19)
        parii.line_spacing=1.1
        parii.bold=True
    
    
    '------------------------------------------'
    '------------------------------------------'
    
# =============================================================================
#  Executive Summary and Methodology slides
# =============================================================================
    "Executive summary slide"
    #when flag=false, the executive summary isnt displayes in the output of the slide but when it's set to true, it is
    flag=True
    if flag:
        summary=datas[9]
        Exec_summ_slide_template =prs.slide_layouts[2]
        
        Exec_summ_slide = prs.slides.add_slide(Exec_summ_slide_template)
       
        a, b, ax, by = Inches(1), Inches(1), Inches(10), Inches(1)
        
        table_shape = Exec_summ_slide.shapes.add_table(3,2 , a, b, ax, by)    #(len(methodology), len(methodology[0]), x, y, cx, cy)
        table_shape.vertical_anchor=MSO_VERTICAL_ANCHOR.BOTTOM
        #specifying IOTA's colour as the colour of the executive summary slide
# =============================================================================
#         
# =============================================================================
        #tbl =  table_shape._element.graphic.graphicData.tbl
        
        #style_id = '{B301B821-A1FF-4177-AEE7-76D212191A09}'
        #tbl[0][-1].text = style_id
# =============================================================================
#         
# =============================================================================
        table1 = table_shape.table
        table1.vertical_anchor=MSO_VERTICAL_ANCHOR.BOTTOM
        #remove table headers
        table1.first_row = False
        table1.columns[0].width=Inches(2.0)

        
        #adding textframe to the title of the executive summary slide and specifying its specifications/dimensions
        Exec_summ_title=Exec_summ_slide.placeholders[0].text_frame.add_paragraph()
        Exec_summ_title.font.color.rgb=RGBColor(27, 113, 171)
        Exec_summ_title.font.size=Pt(25)
        Exec_summ_title.text="Executive Summary"
        
        Exec_summ_slide.placeholders[0].top = Cm(0.6)
        Exec_summ_slide.placeholders[0].left = Cm(2.3)
        Exec_summ_slide.placeholders[0].width = Cm(30)
        Exec_summ_slide.placeholders[0].height = Cm(1.5)
        
        
        
        #specifying the headers of the executive summary slide
        
        cell1 = table1.cell(0,0)
        cell1.text=  "Topic" #methodology[0][0]
        cell1.text_frame.paragraphs[0].alignment=PP_ALIGN.CENTER
        cell2 = table1.cell(0,1)
        cell2.text= "Insights" #methodology[0][1]
        #for i in range(1,len(methodology)):
            #cell3 = table1.cell(i,0)
            #cell3.text=methodology[i][0]
            #cell3.text_frame.paragraphs[0].alignment=PP_ALIGN.CENTER
            #cell4 = table1.cell(i,1)
            #cell4.text=methodology[i][1]
            
            
        #removing other placeholders
        for placeholder in Exec_summ_slide.shapes.placeholders:
            if placeholder.has_text_frame and placeholder.text_frame.text=="":
                twist = placeholder._sp
                twist.getparent().remove(twist)
    else:
        flag=False
        pass
        

    """
    In this section the methodology template is created. We also set up the dimensions of the table that
    is going to contain the information. This is taken from the methodology array which is the one in what we asigned
    the data[7] parameter that contains the array of methodology information. For example:
    methodology =#[["Topic","Insights"],["Affiliate","sense of humor"],["Disclosure","no man sky"]] 
    thus wil create a table like:

    ----------------------------------------------|---------------------------------------
    Topic                                         | Insights
    ----------------------------------------------|---------------------------------------
    Affiliate                                     | sense of humor
    ----------------------------------------------|---------------------------------------
    Disclosure                                    | no man's sky
    ----------------------------------------------|----------------------------------------


    The size of each cell is automatically adapted to the text.
    """
    
    # methodology =[["Topic","Insights"],["Affiliate","sense of humor"],["Disclosure","no man sky"]] #data[7]
    methodology = datas[7]
    methodology_slide_template =prs.slide_layouts[2]
    for box in [2,1,2]:
        
        methodology_slide_contentplaceholder_box = methodology_slide_template.shapes[box]._sp
        methodology_slide_contentplaceholder_box.getparent().remove(methodology_slide_contentplaceholder_box) 



    
  
    methodology_slide = prs.slides.add_slide(methodology_slide_template)
   
    x, y, cx, cy = Inches(1), Inches(1), Inches(18.6), Inches(3)
    #shape = methodology_slide.shapes.add_table(len(methodology), len(methodology[0]), x, y, cx, cy)
    shape = methodology_slide.shapes.add_table(len(methodology), 2, x, y, cx, cy)
    
   
    shape.vertical_anchor=MSO_VERTICAL_ANCHOR.BOTTOM
    #removing header row for methodology table
    table = shape.table
    
    
    
    table.vertical_anchor=MSO_VERTICAL_ANCHOR.BOTTOM
    table.first_row = False
    #table.columns[0].alignment=PP_ALIGN.CENTER
    table.columns[0].width=Inches(2.0)
    #table.columns[1].width=Inches(9)
    
    
    #specifying the dimensions, specifications and name of the methodology slide title
    meth_title=methodology_slide.placeholders[0].text_frame.add_paragraph()
    meth_title.font.color.rgb=RGBColor(27, 113, 171)
    meth_title.font.size=Pt(25)
    meth_title.text="Methodology"
    
    methodology_slide.placeholders[0].top = Cm(0.6)
    methodology_slide.placeholders[0].left = Cm(2.3)
    methodology_slide.placeholders[0].width = Cm(30)
    methodology_slide.placeholders[0].height = Cm(1.5)
    
    #the cells in the methodology slide and their contents
    
    cell1 = table.cell(0,0)
    cell1.text=methodology[0][0]
    cell1.text_frame.paragraphs[0].alignment=PP_ALIGN.CENTER
    
    
    cell2 = table.cell(0,1)
    #cell2.has_text_frame=True
    cell2=cell2.text_frame.paragraphs[0]
    #cell2=cell2.add_run()
    cell2.level=0
    #cell2.text='•  ' + methodology[0][1]
    cell2.text=methodology[0][1]
    
    for i in range(1,len(methodology)):
        cell3 = table.cell(i,0)
        cell3.text=methodology[i][0]
        cell3.text_frame.paragraphs[0].alignment=PP_ALIGN.CENTER
        
        #for k in methodology[i][1:]:
        cell4 = table.cell(i,1)
        cell4=cell4.text_frame.paragraphs[0]
        cell4.level=0
        cell4.text=methodology[i][1]
    
    for placeholder in methodology_slide.shapes.placeholders:
       if placeholder.has_text_frame and placeholder.text_frame.text=="":
           twist = placeholder._sp
           twist.getparent().remove(twist)   

# =============================================================================    
#     #INSIGHTS
# =============================================================================
    """
    Insights_slide = prs.slides.add_slide(methodology_slide_template)
    x, y, cx, cy = Inches(1), Inches(1), Inches(10), Inches(1)
    shape = Insights_slide.shapes.add_table(len(insights), len(insights[0]), x, y, cx, cy)
    table = shape.table
    Insights_slide.placeholders[0].text="Methodology"
    Insights_slide.placeholders[13].text = ""
    cell1 = table.cell(0,0)
    cell1.text=insights[0][0]
    cell2 = table.cell(0,1)
    cell2.text=insights[0][1]
    for i in range(1,len(insights)):
        cell3 = table.cell(i,0)
        cell3.text=insights[i][0]
        cell4 = table.cell(i,1)
        cell4.text=insights[i][1]

    """

# =============================================================================    
#     #GRAPH SLIDES
# =============================================================================
    slide_text = datas[8] 
    title_ppt = datas[1]
    graph_ppt = datas[2]
    source = datas[3]
       
    graph_slide_template = prs.slide_layouts[2]

    #Font shape location
    graph_slide_template.shapes[2].top = Cm(17)
    graph_slide_template.shapes[2].left = Cm(2)  
    graph_slide_template.shapes[2].width = Cm(25)
    graph_slide_template.shapes[2].height = Cm(1)
    
    # Font for source could not be changed, was changed after creation within the iteration    
    graph_slide_template.shapes[2].text_frame.paragraphs[0].font.size = Pt(5)
    graph_slide_template.placeholders[1].text_frame.paragraphs[0].runs[0].font.size = Pt(5)
    graph_slide_template.shapes[2].text_frame.paragraphs[0].runs[0].font.size = Pt(5)
    
    #The graph title's MAGIC happens below
    #print(len(title_ppt))
    #print(len(slide_text))

    for i in range(0,len(title_ppt)):
        graph_slide = prs.slides.add_slide(graph_slide_template)  
        
        graph_ti=graph_slide.placeholders[0].text_frame
        #below sets the alignment of the graph titles to top(BOTTOM is the correct thing. DO NOT CHANGE)
        graph_ti.vertical_anchor=MSO_ANCHOR.BOTTOM
        '-----------------------------------------------------------------------------------------------'
        
        graph_title=graph_ti.add_paragraph()
        graph_title.font.color.rgb=RGBColor(27, 113, 171)
        graph_title.font.size=Pt(25)
        graph_title.text = title_ppt[i]
        graph_slide.placeholders[0].top = Cm(0.3)
        graph_slide.placeholders[0].left = Cm(1.8)
        graph_slide.placeholders[0].width = Cm(30.5054)
        graph_slide.placeholders[0].height = Cm(2.616)
        graph_slide.placeholders[13].text = source
        
        #graph pictures are added...(i think)
        im = Image.open(graph_ppt[i])
        real_width, real_height = im.size
        relation_image = real_width/real_height
        #Specifiesthe dimension of the graphs***ALL OF THEM***
        graph_slide.shapes.add_picture(graph_ppt[i], 
                                       left = Inches(0.8), 
                                       top = Inches(1.90), 
                                       height = Inches(4.80), 
                                       width = 
                                       Inches(5.5*relation_image))
        
        graph_slide.shapes[1].text_frame.paragraphs[0].font.size = Pt(10)
        graph_slide.shapes[2].line.color.rgb = 	RGBColor(0,0,0)
        graph_slide.shapes[2].line.width = Pt(0.25)
        
        
        
            
        sl_t=graph_slide.placeholders[14]
        # sp=sl_tx._sp
        # sp.getparent().remove(sp)
        #sl_tx.line.color=RGBColor(255,255,255)
        line =sl_t.line
        line.color.rgb = RGBColor(255, 255, 255)
        sl_tx=sl_t.text_frame
        #adding a text_frame to the box
        slide_tx=sl_tx.paragraphs[0]
        slide_tx.font.color.rgb=RGBColor(0, 0, 0)
        slide_tx.font.size=Pt(14)
        slide_tx.text = slide_text[i]#sfdgdggg

        
        graph_slide.placeholders[14].top = Inches(1.15)
        graph_slide.placeholders[14].left = Inches(0.71)
        graph_slide.placeholders[14].width = Inches(11.8)
        graph_slide.placeholders[14].height = Inches(0.69)
        
        
        #for shape in graph_slide.placeholders:
    #    print('%d %s' % (shape.placeholder_format.idx, shape.name))
        
        
        #this is the section that adjuste the position of the graphs but creates a bug in the watermark.
        
        '''
        graph_slide.shapes.add_picture(graph_ppt[i],
                                       left = Inches(0.8),
                                       top = Inches(1.5),
                                       height = Inches(5),
                                       width =
                                       Inches(5*relation_image))
        
        graph_slide.shapes[1].text_frame.paragraphs[0].font.size = Pt(10)
        graph_slide_template.shapes[1].top = Cm(17)
        graph_slide_template.shapes[1].left = Cm(5)
        graph_slide.shapes[2].line.color.rgb = 	RGBColor(0,0,0)
        graph_slide.shapes[2].line.width = Pt(0.25)
        '''


# =============================================================================
# Agenda
# =============================================================================
    agenda_slide_template = prs.slide_layouts[4]
    agenda_ppt = datas[0]
    fake_agenda_ppt = datas[4]
    space_occupers = datas[6]
    #print(datas)



    
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)  
    
    #specific level agenda is "1" whic is the level to which the hyperlinks are added
    specific_level_agenda = datas[5] #añade los hyperlinks al elemento de nivel n
   
    
   
    

    def get_slide_jumps(f_agenda,level,initial):
        """
        This function calculates the number of slides needed to each hyperlink based on the fake agenda array

        :params f_agenda: the fake agenda
        :params level: the level on which we are going to set up the hyperlinks on the table of content
        :params initial: this is calculated len(agenda_pg)+EXTRA_CONTENT
        """
        jumps=[]
        slides=[]
        accum=initial
        for i in range(0,len(f_agenda)):
            if(f_agenda[i] == level):
                jumps.append(i)
        slides.append(accum)

        for i in range(0,len(jumps)-1):
            rec = jumps[i+1]-jumps[i]
          
            if(rec == 2 or rec == 1 or rec ==3):
                accum=accum+(rec)
                slides.append(accum)
            else:
                accum=accum+(rec-1)
                slides.append(accum)

    
        return slides
   

    def level(number_agenda, level, body,to_slide):
            '''add the text and level of a paragraph'''
            shape = body
            nivel = body.text_frame.add_paragraph()
            nivel.clear()
            nivel.level = level
            #ivel.text = agenda_ppt[first_item:last_item][number_agenda].split(".")[1]
            nivel.font.color.rgb = RGBColor(0, 0, 0)
           
            nivel = nivel.add_run()
            nivel.font.size=Pt(15)
            nivel.text = agenda_ppt[first_item:last_item][number_agenda].split(".")[1]
            
            
            if level == specific_level_agenda:             
                shape.click_action.target_slide = prs.slides[to_slide]
                nivel.hyperlink.address = shape.click_action.hyperlink.address
                nivel.hyperlink._hlinkClick.action = shape.click_action.hyperlink._hlink.action
                nivel.hyperlink._hlinkClick.rId = shape.click_action.hyperlink._hlink.rId
                shape.click_action.target_slide = None
            
            return nivel
    def get_agenda_values(agenda):
     """
     this function gets the first character of the list, in this case it generates arrays like ['0','1','2','1'] 

     :param agenda: list that contains all the information used for generate the table of content.
     """
     values=[]
     for i in range(0,len(agenda)):
         values.append(agenda[i][0])
     return values
     #this function appends all the first characters of agenda (the digits 0 or 1 that starts each agenda) to values,
     # and has a length equal to the length of agenda. and since that's the case, each item in values is seen as string!
 

    agenda_values=[int(i) for i in get_agenda_values(agenda_ppt)]
    fake_agenda_values=[int(i) for i in get_agenda_values(fake_agenda_ppt)]
    #these two variables agenda_values and fake_agenda_values turn the values in 'value' into integers
    #so they can be used to get the levels of the items in agenda_ppt and fake_agenda_ppt


    #THIS WHOLE SECTION IS FOR ELIMINATE EVERY VALUE ON THE LIST THAT FOR THE CURRENT MOMENT DOESN'T HOLD ANY
    #SLIDE NOR BELONG TO THE LEVEL, SO IT IS UNNECESARRY FOR THE HYPERLINK NUMBER OF JUMPS CALCULATION.
    temp=[]
    slides=[]
    for j in range(0, len(fake_agenda_values)):
        if(fake_agenda_values[j] in space_occupers):
            temp.append(fake_agenda_values[j])
        elif(fake_agenda_values[j] == specific_level_agenda):
            temp.append(fake_agenda_values[j])

    
    #"""def split(arr, size):
      #   agenda_iii = []
      #   while len(arr) > size:
     #        pice = arr[:size]
     #        agenda_iii.append(pice)
     #        arr   = arr[size:]
        # agenda_iii.append(arr)
    #     return agenda_iii"""
    #def get_step():
        #listA = [11, 18, 19, 21, 29, 46]
        
    #splits = np.array_split(agenda_ppt, last_item)
    
    #for arrays in splits:
        #print(list(array))
        
        
        
    #print(agenda_ppt)    
    #declaring sekunda as am empty list
    sekunda=[]
    #loop to get number of lines in each item in "agenda_ppt" as integerss
    for m in range(0, len(agenda_ppt)):
        if len(agenda_ppt[m])>=62<=124:
            pipl=2
        elif len(agenda_ppt[m])>=124<=174:
            pipl=3
        elif len(agenda_ppt[m])>=174:
            pipl=4
        else:
            pipl=1
        sekunda.append(pipl)
    #sekunda is a list that contains the number of lines each agenda index has
    
            
    # Cantidad de items que caben comodamente en una hoja
   
    """
    made_in_kafanchan is a function that examines the list 'sekunda', takes in a starting point 'First' and
    adds up each consequtive number in the list while returning the index of the number that sums up to 17. It then adds
    1 to that index and returns it as step and also repeats the process over again until the list runs out
    
    """
    def made_in_kafanchan(First, sekunda):
        sum = 0
        step=0
        for makay in range(First, len(sekunda)):
            sum+=sekunda[makay]
            if sum>17:
                break
            step=makay+1
        return step
    #print(sekunda)
    first_item = 0
    #the last item will be a function of the first item since the former is not defined and varies according to the 
    #number of sentences each group of items in one column has
    last_item= made_in_kafanchan(first_item, sekunda)
    #agenda pg will also depend on the new list sekunda because it has the total number of sentences agenda_ppt forms
    agenda_pg = list(range(0, int(math.ceil(sum(sekunda)/(2*(16))))))
    #     #MODIFIED BY OSCAR
   # print(agenda_pg)
    fix_extra_content = 0
    for fix in agenda_pg:
        fix_extra_content = 1 + fix_extra_content
    # =============================================================================
    EXTRA_CONTENT = 6 - fix_extra_content 
   
    slides=get_slide_jumps(temp,specific_level_agenda,len(agenda_pg)+EXTRA_CONTENT)
    
    
    
    #there are two cases that work differently because of 'pg'. when the length of agenda_pg is more than 2 and when it is 1
    #when it's len is just 1, it is 0
    #since agenda_pg is a list of consecutive numbers, the highest number will give the last page
    if len(agenda_pg)>=2:
    	consi=agenda_pg[-1]
    else:
        consi=agenda_pg[0]
    #consi==1 makes a special case for when there are two agenda pages
    if consi>= 2 or consi ==1: 
        
        #adding the agenda pages for cases where the number of pages are 2 or more
        for pg in agenda_pg:
            agenda_slide = prs.slides.add_slide(agenda_slide_template) 
            '-----------------from here----------------------------------------'
            ####################################################################
            agenda_slide.shapes.placeholders[13].top=Inches(0.9)
            agenda_slide.shapes.placeholders[13].left=Inches(0.6)
            agenda_slide.shapes.placeholders[13].width=Inches(8)
            agenda_slide.shapes.placeholders[13].height=Inches(0.4)
            
            click_nav=agenda_slide.shapes.placeholders[13].text_frame.paragraphs[0]#.text = "Click to navigate"
            click_nav.clear()
            click=click_nav.add_run()
            click.text="Hold Ctrl and Click to navigate"
            click.font.color.rgb = RGBColor(27, 113, 171)
            click.font.name="Mulish"
            
            agenda_slide.shapes.placeholders[0].left=Inches(0.6)
            agenda_slide.shapes.placeholders[0].top=Inches(0.5)
            agenda_slide.shapes.placeholders[0].width=Inches(11)
            agenda_slide.shapes.placeholders[0].height=Inches(0.5)
            
            agenda_slide.shapes.placeholders[1].top=Inches(0.86)
            agenda_slide.shapes.placeholders[1].left=Inches(0.6)
            agenda_slide.shapes.placeholders[1].width=Inches(6.1)
            agenda_slide.shapes.placeholders[1].height=Inches(5.91)
        
            
            agenda_slide.shapes.placeholders[14].top=Inches(0.86)
            agenda_slide.shapes.placeholders[14].left=Inches(6.7)
            agenda_slide.shapes.placeholders[14].width=Inches(6.1)
            agenda_slide.shapes.placeholders[14].height=Inches(5.91)     
            #######################################################################
            titv=agenda_slide.shapes.placeholders[0].text_frame.paragraphs[0]#.text = "Click to navigate"
            titv.clear()
            ick=titv.add_run()
            ick.text="Index (pg. %d/%s)" % (pg+1, agenda_pg[-1]+1)
            ick.font.color.rgb = RGBColor(27, 113, 171)
            ick.font.name="Mulish"
            ick.font.size=Pt(25)
            ######################################################################
            '----------------------to here-------------------------------------------------'
            'specifies the contents that apply to each pages in the agenda slides including the specifications'
            ' and dimentions, text and text sizes of the titles for the case were agenda pages are more than 1'
            'and below specifies the cntents of the left and right placeholders, encompassing most of the functions'
            'before this loop including level where the MAGIC of the agenda slides happen'

            for lemi in [1,14]:     
               
                body = agenda_slide.shapes.placeholders[lemi]
                #body.shapes[1].text_frame.paragraphs.font.rgb = RGBColor(0, 0, 255)
                cont = 0

                for numero_agenda in range(0,len(agenda_ppt[first_item:last_item])):
                    
                    level_agenda = agenda_values[first_item:last_item][numero_agenda]
                    #print(first_item, last_item, numero_agenda)
                    level(numero_agenda, level_agenda, body, slides[cont])
                    #umero_agenda.font.color.rgb = RGBColor(0, 0, 0)
                    if level_agenda == specific_level_agenda and cont < (len(slides)-1):
                        
                        cont=cont+1
                #step will differ for each column in each page and will depend on the number of sentences from agenda_ppt 
                #each column can contain and will follow from the last item in the previous column
                step = made_in_kafanchan(last_item, sekunda)
                #the first item of a colum will follow from the last item of the previous column
                first_item = last_item
                last_item = step 
                'and up till this point'
            #below removes all excess placeholders in the agenda slides for when the pages are more than one    
            for placeholder in agenda_slide.shapes.placeholders:
               if placeholder.has_text_frame and placeholder.text_frame.text=="":
                   twist = placeholder._sp
                   twist.getparent().remove(twist)
  
    else:
        #when the number of agenda pages is 1, pg will be the first item in agenda_pg i.e '0'
        pg = agenda_pg[0]
        
        
        #*********FROM HERE*******
        agenda_slide = prs.slides.add_slide(agenda_slide_template)
        agenda_slide.shapes.placeholders[13].top=Inches(0.9)
        agenda_slide.shapes.placeholders[13].left=Inches(0.6)
        agenda_slide.shapes.placeholders[13].width=Inches(8)
        agenda_slide.shapes.placeholders[13].height=Inches(0.4)
        
        click_nav=agenda_slide.shapes.placeholders[13].text_frame.paragraphs[0]#.text = "Click to navigate"
        click_nav.clear()
        click=click_nav.add_run()
        click.text="Hold Ctrl and Click to navigate"
        click.font.color.rgb = RGBColor(27, 113, 171)
        click.font.name="Mulish"
        
        agenda_slide.shapes.placeholders[0].left=Inches(0.6)
        agenda_slide.shapes.placeholders[0].top=Inches(0.5)
        agenda_slide.shapes.placeholders[0].width=Inches(11)
        agenda_slide.shapes.placeholders[0].height=Inches(0.5)
        
        agenda_slide.shapes.placeholders[1].top=Inches(0.86)
        agenda_slide.shapes.placeholders[1].left=Inches(0.6)
        agenda_slide.shapes.placeholders[1].width=Inches(6.1)
        agenda_slide.shapes.placeholders[1].height=Inches(5.91)
    
        
        agenda_slide.shapes.placeholders[14].top=Inches(0.86)
        agenda_slide.shapes.placeholders[14].left=Inches(6.7)
        agenda_slide.shapes.placeholders[14].width=Inches(6.1)
        agenda_slide.shapes.placeholders[14].height=Inches(5.91) 
        #######################################################################
        titv=agenda_slide.shapes.placeholders[0].text_frame.paragraphs[0]#.text = "Click to navigate"
        titv.clear()
        ick=titv.add_run()
        ick.text="Index (pg. %d/%s)" % (pg+1, agenda_pg[0]+1)
        ick.font.color.rgb = RGBColor(27, 113, 171)
        ick.font.name="Mulish"
        ick.font.size=Pt(25)
        ########################**********TO HERE*************##############################################
        #specifies the dimensions, specifications, title texts and font for when the number of pages of index is 1
        
        #SPECIFIES THE CONTENTS IN EACH COLUMN OF THE PAGE including the text and levels and functions defined before this point
        for lemi in [1,14]:     
           
            body = agenda_slide.shapes.placeholders[lemi]
            #body.shapes[1].text_frame.paragraphs.font.rgb = RGBColor(0, 0, 255)
            cont = 0
            for numero_agenda in range(0,len(agenda_ppt[first_item:last_item])):        
                level_agenda = agenda_values[first_item:last_item][numero_agenda]
    
                level(numero_agenda, level_agenda, body, slides[cont])
                #umero_agenda.font.color.rgb = RGBColor(0, 0, 0)
                if level_agenda == specific_level_agenda and cont < (len(slides)-1):
                    
                    cont+=1
            
            step = made_in_kafanchan(last_item, sekunda)
            first_item = last_item
            last_item = step 
            #UP TO THIS POINT
    #bElow then removes any excess placeholder that contains no text      
    for placeholder in agenda_slide.shapes.placeholders:
       if placeholder.has_text_frame and placeholder.text_frame.text=="":
           twist = placeholder._sp
           twist.getparent().remove(twist)

    
# =============================================================================
# keywords and data dictionary
# =============================================================================
    
    year = int(re.search('\d{4}', datas[3]).group(0))
    
    source = "0.Source: Year %d*" % year    
    

    def transform_multindex_col(data_dictionary):
        ''' function to transform multindex columns to normal keys'''
        final_keys = []
        last_keys = []
        counter = 0
        for key in data_dictionary.keys():
            final_keys.append([sub_key for sub_key in key if sub_key[0:7] != "Unnamed"])
        for final_key in final_keys:
            last_key = " ".join(final_key)
            last_keys.append(last_key.capitalize())    
        
        for key in list(data_dictionary.keys()):        
            data_dictionary[last_keys[counter]] = data_dictionary.pop(key)
            counter = counter + 1
        return data_dictionary
    
    if data_dictionary == {}:
        variables_texto = ""
    else:
        if type(list(data_dictionary.keys())[0]) is tuple:
            data_dictionary = transform_multindex_col(data_dictionary) 
                    
        variables_texto = ""
        for key in data_dictionary.keys():
            try:
                variables_texto=variables_texto+"*1."+key+"*2."+", ".join(data_dictionary[key])
            except:
                variables_texto = variables_texto

        variables_texto = "*0.Categorical variables"+variables_texto
     
    number_of_characters = 800

    #total string for agenda dictionary
    #removed  keywords+answers+
    text = source + variables_texto
    
    len(text)
    
     #total item list 
    text.split("*")
    len(text.split("*"))
    
    # calculating amount of required pages
    
    # print('metadata_pg', metadata_pg)
    text_amount = []
    len_acum_text = 0
    for item in text.split("*"):
        len_acum_text = len_acum_text+ len(item) 
        text_amount.append(len_acum_text)
    
    # DATAFRAME, where there are the items with their accumulative length
    skip_page_df = pd.concat([pd.Series(text.split("*"), name = "item"),
                  pd.Series(text_amount, name = "text_amount")], axis = 1)
    


    def level_meta_data(item, body, items_list):
        '''add the text and level of a paragraph'''
        # if the name of item matches with the first name of the list do:

        if (len(item) > 0):
            if item.split(".")[1] == items_list[0].split(".")[1]:
                if int(item.split(".")[0]) == 0:
                    # if the level of the first item is equal to 0 treat it as level 0,
                    # without adding paragraph
                    line = body.text_frame
                    line.text = item.split(".")[1]
                    line.level = int(item.split(".")[0])
                else:
                    # else treat it as a inferior level
                    line = body.text_frame.add_paragraph()
                    line.text = item.split(".")[1]
                    line.level = int(item.split(".")[0])
            else:
                #if is not the item is not the first item of the list do this:
                line = body.text_frame.add_paragraph()
                line.text = item.split(".")[1]
                line.level = int(item.split(".")[0])
 
    '''
    if len(answers.split("*")) >= 6:
        fixed = 1
    else:# NO FIXING
        fixed = 0
    '''   
    top_index = 0
    bottom_index = 0
        
    skip_page_df.item.iloc[0:12].values
    skip_page_df.item.iloc[12:21].values
    
    last_item = number_of_characters
    step = number_of_characters
    metadata_pg = list(range(0, (int(math.ceil(len(text)/(2*number_of_characters))))))
    
    for pg in metadata_pg:   
        # for each metadata requeried pages do this:
        meta_slide_template=prs.slide_layouts[5]
        
        meta_data_slide = prs.slides.add_slide(meta_slide_template) 
        
        meta=meta_data_slide.shapes.placeholders[0].text_frame.paragraphs[0]
        #meta.clear()
        meta.font.color.rgb = RGBColor(27, 113, 171)
        meta.font.name="Mulish"
        meta.font.size=Pt(25)
        meta.text = "More Information (pg. %d/%d)" % (pg+1, metadata_pg[-1]+1)
        
        #fixing box location

        meta_data_slide.shapes[1].top = Inches(1.04)
        meta_data_slide.shapes[1].left = Inches(0.5)  
        meta_data_slide.shapes[1].width = Inches(6.08)
        meta_data_slide.shapes[1].height = Inches(5.7)
        
        meta_data_slide.shapes[2].top = Inches(1.04)
        meta_data_slide.shapes[2].left = Inches(6.58)  
        meta_data_slide.shapes[2].width = Inches(6.08)
        meta_data_slide.shapes[2].height = Inches(5.7)
        
        
        
        meta_data_slide.shapes.title.top = Inches(0.5)
        meta_data_slide.shapes.title.left = Inches(0.5)
        meta_data_slide.shapes.title.width = Inches(12.16)
        meta_data_slide.shapes.title.height = Inches(0.52)
        visit=meta_data_slide.shapes.placeholders[11].text_frame.paragraphs[0]
        visa=visit.add_run()
        visa.text="Visit https://www.iotaq.tech/"
        visa.font.color.rgb = RGBColor(27, 113, 171)
        visa.font.name="Mulish"
        visa.font.size=Pt(18)
        hlink=visa.hyperlink
        hlink.address='https://www.iotaq.tech/'
        for uni in [1,2]:
            body = meta_data_slide.shapes.placeholders[uni]

            close_value = int(max(list(skip_page_df["text_amount"].apply(lambda x: (x-last_item)).apply(lambda x: x if x < 0
                                                                                 else float("-inf"))))+last_item) 
            '*'
            top_index = skip_page_df[skip_page_df["text_amount"] == close_value].index[0]+1 # add 1 because when using iloc it will take dataset lenght      
            items_list = list(skip_page_df.item.iloc[bottom_index:(top_index)].values) #+fixed variable 
            
            for item in items_list:
                # for each item of the item list belonging to each page do this
                level_meta_data(item, body, items_list)
                
            # for shape in meta_data_slide.shapes:
            #     for paragraph in meta_data_slide.shapes[1].text_frame.paragraphs:
            #         paragraph.font.size = Pt(14)
            
            
            
            bottom_index = top_index  
            last_item = last_item + step
        
        for placeholder in meta_data_slide.shapes.placeholders:
            if placeholder.has_text_frame and placeholder.text_frame.text=="":
                twist = placeholder._sp
                twist.getparent().remove(twist)
        

# # =============================================================================
# Delete template
# =============================================================================
    def delete_slide(self, presentation,  index):
            xml_slides = presentation.slides._sldIdLst  # pylint: disable=W0212
            slides = list(xml_slides)
            xml_slides.remove(slides[index])
            
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)  
    for j in range(0,3):
        xml_slides.remove(slides[j]) 
    
        

# =============================================================================
# Rearange slides (not possible)
# =============================================================================



    #REPOSITIONING THE SLIDES ALSO DEPENDS ON WHETHER THEERE IS JUST 1 AGENDA PAGE OR THERE ARE TWO OR MORE PAGES    
    def move(self, presentation, old_index, new_index):
        xml_slides = presentation.slides._sldIdLst  # pylint: disable=W0212
        slides = list(xml_slides)
        xml_slides.remove(slides[old_index])
        xml_slides.insert(new_index, slides[old_index])
    #If there are two or more pages, the total number of pages-number of agenda pages-number of meta_data pages will
    #give the starting point i.e. index of where agenda pages were formed and the total number of pages-meta_data pages
    #will give the end point WHILE "Position" will be determined by the integer on the R.H.S which specifies where the 
    #counter for where the agenda will be moved to will begin.
    counter = 0
    if consi>=2 or consi==1:
        agenda_old_indexes = list(range(len(prs.slides)-len(metadata_pg)-len(agenda_pg),len(prs.slides)-len(metadata_pg)))
        #agenda_old_indexes.reverse()
        for index in agenda_old_indexes:
            position = 1 + counter
            move(prs.slides[index], prs, index, position)
            counter = counter + 1
    #AND when the number of pages of the agenda is 1, the page will be formed at the index "total number of slides" - 
    #                                                                                   "number of meta_data slides" - 
    #                                                                           "1" which is the length od string(consi)
    #as in the ELSE statement BELOW
    else:
        agenda_old_indexes = len(prs.slides)-len(metadata_pg)-len(str(consi))
        index = agenda_old_indexes
        position = 1 + counter
        move(prs.slides[index], prs, index, position)
    #BELOW forms the placeholders(FOOTERS) for the PAGE NUMBERS in this slide.
    xml_slides = prs.slides._sldIdLst
    count=0
    
    
    for slide in prs.slides:
        
        if watermark==True:
                left=Inches(0)
                top=Inches(0)
                height=prs.slide_height
                width=prs.slide_width
                pic=slide.shapes.add_picture(water_mark, left, top, width, height)
                # slide.shapes._spTree.remove(pic._element)
                # slide.shapes._spTree.insert(2, pic._element)
        else:
            pass

        
        #BELOW IF STATEMENT makes the counter start from index 1 and not 0 thereby aloowing the title page not to be numbered
        if count>0:
            
            #BELOW specifies the size of the Textbox that contains the page numbers
            top=Inches(6.8)
            left = Inches(12.2)
            width = height = Inches(1)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            #adding a text_frame to the box
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            #specifying the colour of the text in the box
            run.font.color.rgb = RGBColor(0, 128, 255)
            #String for the text in box in each slide
            run.text = str(count+1) + '/' + str(len(xml_slides))
        #ELSE statement specifying what happens to the pages not numbered i.e THE TITLE SLIDE
        else:
            
            pass
        #COUNTER
        count+=1
        #print('::::::::::::::::::::::::')
                #print(level_agenda)
    
    prs.save('reports\\{}.pptx'.format(name))   
#    PPTtoPDF('reports\\{}.pptx'.format(name), 'reports\\{}'.format(name))
    print("pptx and pdf files have been created successfully in folder reports")
    
def PPTtoPDF(inputFileName, outputFileName, formatType = 32):
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1

    inputFileName = os.getcwd()+"\\"+inputFileName
    outputFileName =os.getcwd()+"\\"+outputFileName
    if outputFileName[-3:] != 'pdf':
        outputFileName = outputFileName + ".pdf"
    deck = powerpoint.Presentations.Open(inputFileName)
    deck.SaveAs(outputFileName, formatType) # formatType = 32 for ppt to pdf
    deck.Close()
    powerpoint.Quit()










    